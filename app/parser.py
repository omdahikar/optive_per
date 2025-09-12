import fitz
import pytesseract
from PIL import Image
import openpyxl
import os
from pptx import Presentation
from rich.console import Console

console = Console()

# --- ⚠️ CRITICAL TESSERACT CONFIGURATION ⚠️ ---
TESSERACT_CMD_PATH = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
# ----------------------------------------------------

try:
    if TESSERACT_CMD_PATH and os.path.exists(TESSERACT_CMD_PATH):
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD_PATH
    else:
        pytesseract.get_tesseract_version()
except Exception:
    console.print("[bold red]FATAL TESSERACT ERROR:[/bold red] 'tesseract.exe' was not found.")
    console.print("  1. Make sure Tesseract is installed.")
    console.print("  2. [bold]Open 'app/parser.py' and set the TESSERACT_CMD_PATH variable.[/bold]")
    exit()

def extract_content(file_path):
    if not os.path.exists(file_path): return None
    file_extension = os.path.splitext(file_path)[1].lower()
    base_name = os.path.basename(file_path)
    console.print(f"  📄 [bold]Parsing:[/] Reading content from [cyan]{base_name}[/cyan]...")
    extracted_text = ""
    try:
        if file_extension == '.pdf':
            with fitz.open(file_path) as doc:
                extracted_text = "".join(page.get_text() for page in doc)
        elif file_extension in ['.png', '.jpeg', '.jpg']:
            # Better OCR with configuration and feedback
            custom_config = r'--oem 3 --psm 6'
            extracted_text = pytesseract.image_to_string(Image.open(file_path), config=custom_config)
            if not extracted_text.strip():
                console.print("[yellow]  ↪ Tesseract Warning: OCR could not find any text in this image. This can happen with low-resolution or blurry images.[/yellow]")
        elif file_extension == '.xlsx':
            workbook = openpyxl.load_workbook(file_path)
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    extracted_text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
        else:
            return None
    except Exception as e:
        console.print(f"[bold red]  ↪ Error parsing {base_name}:[/bold red] {e}")
        return None
    return {"text": extracted_text.strip(), "images": []}