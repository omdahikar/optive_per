from pptx import Presentation
from pptx.util import Pt, Inches
from rich.console import Console
import os

console = Console()

def generate_report(analysis_results, template_path, output_dir):
    """
    Generates the final report by replacing a placeholder shape on the template slide
    with a new, perfectly formatted table. This is the most robust method.
    """
    console.print("\n[bold]📊 Generating Final PowerPoint Report...[/bold]")
    try:
        prs = Presentation(template_path)
    except Exception as e:
        console.print(f"[bold red]  ↪ Error opening template '{template_path}': {e}[/bold red]")
        return

    # --- FIND THE SLIDE AND PLACEHOLDER ---
    target_slide = None
    placeholder = None
    for slide in prs.slides:
        for shape in slide.shapes:
            # We use the shape's default name 'Content Placeholder 2' as a reliable marker
            if shape.name == 'Content Placeholder 2':
                target_slide = slide
                placeholder = shape
                break
        if placeholder:
            break
            
    if not placeholder:
        console.print("[bold red]  ↪ Error: Could not find the 'Content Placeholder 2' on the 'File Analysis' slide.[/bold red]")
        console.print("[bold red]  ↪ Please ensure the template slide has a main content body.[/bold red]")
        return
    
    # --- CREATE THE NEW TABLE ---
    left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
    
    # Remove the old placeholder shape
    sp = placeholder._sp
    target_slide.shapes._spTree.remove(sp)
    
    # Add a new, perfectly sized table in its place
    table = target_slide.shapes.add_table(len(analysis_results) + 1, 4, left, top, width, height).table
    
    # --- FORMAT AND POPULATE TABLE FOR A PROFESSIONAL LOOK ---
    table.columns[0].width = Inches(1.5) # File Name
    table.columns[1].width = Inches(0.8) # File Type
    table.columns[2].width = Inches(3.5) # Description
    table.columns[3].width = Inches(4.0) # Insights

    headers = ["File Name", "File Type", "File Description", "Insights & Extracted Features"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    for idx, result in enumerate(analysis_results):
        row_idx = idx + 1
        table.cell(row_idx, 0).text = result.get('file_name', 'N/A')
        table.cell(row_idx, 1).text = result.get('file_type', 'N/A')
        table.cell(row_idx, 2).text = result.get('description', 'No description generated.')
        table.cell(row_idx, 3).text = "\n".join(result.get('insights', ['- No insights extracted.']))
    
    for row in table.rows:
        for cell in row.cells:
            if not cell.text_frame.paragraphs[0].font.bold: # Don't resize header
                for paragraph in cell.text_frame.paragraphs:
                     paragraph.font.size = Pt(10)

    output_path = os.path.join(output_dir, "Final_AI_Analysis_Report.pptx")
    try:
        prs.save(output_path)
        console.print(f"[bold green]  ↪ Success![/bold green] Report saved to [underline]{output_path}[/underline]")
    except Exception as e:
        console.print(f"[bold red]  ↪ Error saving final report: {e}[/bold red]")